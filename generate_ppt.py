#!/usr/bin/env python3
"""Generate comprehensive project presentation for MGA Tech Consulting."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Constants ---
BG_DARK = RGBColor(0x11, 0x18, 0x27)
BG_SLATE = RGBColor(0x1E, 0x29, 0x3B)
BG_CARD = RGBColor(0x25, 0x30, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
GRAY_MID = RGBColor(0x94, 0xA3, 0xB8)
SKY = RGBColor(0x38, 0xBD, 0xFE)
PURPLE = RGBColor(0xC0, 0x84, 0xFC)
GREEN = RGBColor(0x4A, 0xDE, 0x8A)
AMBER = RGBColor(0xFB, 0xC0, 0x2D)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
INDIGO = RGBColor(0x81, 0x8C, 0xF8)
RED_ACCENT = RGBColor(0xF8, 0x71, 0x71)
FONT_MAIN = "Calibri"
FONT_CODE = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        solid = shape.fill._fill.find(".//a:solidFill", nsmap)
        if solid is not None:
            srgb = solid.find("a:srgbClr", nsmap)
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, f"{{{nsmap['a']}}}alpha")
                alpha_elem.set("val", str(int(alpha * 1000)))
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font=FONT_MAIN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=GRAY_LIGHT, spacing=Pt(6), font=FONT_MAIN, bullet_char="●"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = spacing
    return txBox


def add_card(slide, left, top, width, height, bg_color=BG_CARD, border_color=None):
    shape = add_rect(slide, left, top, width, height, bg_color)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    return shape


def add_section_number(slide, number, left=Inches(0.5), top=Inches(0.4), size=14, color=SKY):
    add_textbox(slide, left, top, Inches(1), Inches(0.4),
                f"0{number}", font_size=size, color=color, bold=True, font=FONT_CODE)


def add_header_bar(slide, slide_num):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SKY)
    add_textbox(slide, Inches(11.8), Inches(6.9), Inches(1.3), Inches(0.4),
                f"{slide_num}", font_size=11, color=GRAY_MID, alignment=PP_ALIGN.RIGHT, font=FONT_CODE)
    add_textbox(slide, Inches(0.4), Inches(6.9), Inches(3), Inches(0.4),
                "www.mgatc.com", font_size=10, color=GRAY_MID, font=FONT_CODE)


# ========================================
# SLIDE 1 - TITLE / COVER
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(7.5), BG_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), SKY)
add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), SKY)
add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), SKY)
add_rect(slide, Inches(13.25), Inches(0), Inches(0.08), Inches(7.5), SKY)

add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.8),
            "MGA TECH CONSULTING", font_size=42, color=SKY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.6),
            "Presencia Digital Integral — Portfolio, Consultoría & Productos", font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.5), Inches(3.1), Inches(4.3), Inches(0.03), SKY)

add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
            "Mariano Gobea Alcoba  |  Data & Analytics Technical Leader  |  Buenos Aires, Argentina",
            font_size=16, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.5),
            "www.mgatc.com  •  mariano@mgatc.com  •  2026",
            font_size=14, color=GRAY_MID, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(3.5), Inches(4.9), Inches(6.3), Inches(1.8), BG_CARD)
stats = [
    "3 sitios web en producción",
    "6 servicios de consultoría",
    "80+ artículos técnicos",
    "6 flujos de automatización n8n",
]
add_bullet_list(slide, Inches(3.7), Inches(5.0), Inches(5.9), Inches(1.6),
                stats, font_size=14, color=GRAY_LIGHT, bullet_char="▸")

# ========================================
# SLIDE 2 - AGENDA
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 2)
add_section_number(slide, 1)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Agenda", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), SKY)

agenda_items = [
    ("01", "Visión General del Proyecto", "Monorepo, arquitectura y stack tecnológico"),
    ("02", "Sitio Principal — MGA Tech Consulting", "Servicios B2B, pricing, conversión y lead gen"),
    ("03", "Portfolio Interactivo & Blog Técnico", "CV digital, 80+ posts, herramientas financieras"),
    ("04", "Sitios de Clientes", "Neil Climatizadores & El Portugués S.A."),
    ("05", "Automatizaciones & Infraestructura", "n8n, Supabase, CI/CD, hosting, analytics"),
    ("06", "SEO, Marketing & Growth", "Estrategia, métricas, oportunidades de mejora"),
    ("07", "Pricing & Modelo de Negocio", "Benchmark, recomendaciones, proyecciones"),
    ("08", "Roadmap & Próximos Pasos", "Hoja de ruta y priorización"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.9) + Inches(i * 0.65)
    add_card(slide, Inches(1), y, Inches(11.3), Inches(0.55), BG_CARD)
    add_textbox(slide, Inches(1.2), y + Inches(0.05), Inches(0.5), Inches(0.45),
                num, font_size=18, color=SKY, bold=True, font=FONT_CODE)
    add_textbox(slide, Inches(1.8), y + Inches(0.02), Inches(5), Inches(0.3),
                title, font_size=16, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.8), y + Inches(0.28), Inches(10), Inches(0.25),
                desc, font_size=12, color=GRAY_MID)

# ========================================
# SLIDE 3 - VISION GENERAL
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 3)
add_section_number(slide, 2)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Visión General del Proyecto", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), SKY)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6),
            "Monorepo que contiene la presencia digital completa de Mariano Gobea Alcoba y MGA Tech Consulting",
            font_size=16, color=GRAY_LIGHT)

cards = [
    ("🏢  Sitio Principal", "www.mgatc.com\nConsultoría + Portfolio + Blog + Recursos", SKY),
    ("🏭  Neil Climatizadores", "www.mgatc.com/neil-site/\nClimatización vehicular — 6 idiomas", PURPLE),
    ("📦  El Portugués S.A.", "www.mgatc.com/elportugues-site/\nLogística y distribución — ISO 9001", TEAL),
]

for i, (title, desc, accent) in enumerate(cards):
    x = Inches(1) + Inches(i * 3.9)
    add_card(slide, x, Inches(2.7), Inches(3.6), Inches(1.8), BG_CARD, accent)
    add_textbox(slide, x + Inches(0.2), Inches(2.85), Inches(3.2), Inches(0.5),
                title, font_size=18, color=accent, bold=True)
    add_textbox(slide, x + Inches(0.2), Inches(3.4), Inches(3.2), Inches(0.9),
                desc, font_size=14, color=GRAY_LIGHT)

add_textbox(slide, Inches(1), Inches(4.9), Inches(11), Inches(0.5),
            "Arquitectura Técnica", font_size=22, color=WHITE, bold=True)

tech_cols = [
    ("Frontend", "Next.js 14 · TypeScript · Tailwind CSS\nFramer Motion · Lucide React", SKY),
    ("Backend / Data", "Supabase (PostgreSQL) · Row Level Security\n24+ tablas · 11 migraciones", GREEN),
    ("Infraestructura", "Cloudflare Pages · GitHub Actions CI/CD\nCloudflare DNS · Email Routing", PURPLE),
    ("Automatización", "n8n (6 workflows) · OpenRouter (LLMs)\nBrevo SMTP · Google Sheets", AMBER),
]

for i, (title, desc, accent) in enumerate(tech_cols):
    x = Inches(1) + Inches(i * 3.0)
    add_card(slide, x, Inches(5.5), Inches(2.8), Inches(1.5), BG_SLATE, accent)
    add_textbox(slide, x + Inches(0.15), Inches(5.55), Inches(2.5), Inches(0.35),
                title, font_size=15, color=accent, bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(5.95), Inches(2.5), Inches(0.9),
                desc, font_size=12, color=GRAY_LIGHT)

# ========================================
# SLIDE 4 - SITIO PRINCIPAL (CONSULTING)
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 4)
add_section_number(slide, 3)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "MGA Tech Consulting — Consultoría B2B", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), SKY)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
            "Orientada a PyMEs argentinas que quieren automatizar, implementar BI o adoptar IA sin equipo técnico propio.",
            font_size=16, color=GRAY_LIGHT)

add_textbox(slide, Inches(1), Inches(2.6), Inches(5), Inches(0.4),
            "Propuesta de Valor", font_size=20, color=SKY, bold=True)

value_props = [
    "Reducción de costos operativos hasta 80%",
    "Primera automatización 100% GRATIS (lead magnet)",
    "Desarrollo ágil con ROI medible en 2-3 meses",
    "Expertise de MercadoLibre aplicado a PyMEs",
]
add_bullet_list(slide, Inches(1), Inches(3.1), Inches(5.5), Inches(2),
                value_props, font_size=14, color=GRAY_LIGHT)

add_textbox(slide, Inches(7), Inches(2.6), Inches(5), Inches(0.4),
            "Componentes de la Página", font_size=20, color=SKY, bold=True)

components = [
    "Hero con parallax + typing animation",
    "Paradigm Shift — propuesta de valor",
    "Free Offer Banner — automatización gratis",
    "6 Servicios detallados con modales",
    "Success Stories + Client Portfolio",
    "Pricing Overview (3 tiers)",
    "Case Studies + Testimonials",
    "Process (4-step methodology)",
    "Contact Form (WhatsApp + Email)",
    "Newsletter Banner — The Data Digest",
    "Floating CTA + Proposal Modal",
]
add_bullet_list(slide, Inches(7), Inches(3.1), Inches(5.5), Inches(3.5),
                components, font_size=13, color=GRAY_LIGHT, bullet_char="▸")

add_textbox(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.4),
            "6 Servicios Ofrecidos", font_size=20, color=SKY, bold=True)

services_data = [
    ("Automatización", "RPA, n8n, Python, webhooks\n$100/mes o $600 proyecto", SKY),
    ("IA Aplicada", "Chatbots WhatsApp, RAG, ML\n$200/mes o $800 proyecto", PURPLE),
    ("BI & Datos", "BigQuery, Looker, dashboards\n$200/mes o $800 proyecto", CYAN),
    ("Mentorías 1-a-1", "Carrera tech, entrevistas, arquitectura\n$100/sesión · $350/pack 4", AMBER),
    ("Cursos Grupales", "Máx 4 personas · SQL, Python, GenAI\n$600-$850 por curso", INDIGO),
    ("Reclutamiento Tech", "Evaluación técnica real por Tech Lead\n$3,000 success fee", TEAL),
]

for i, (title, desc, accent) in enumerate(services_data):
    x = Inches(1) + Inches(i * 2.05)
    add_card(slide, x, Inches(5.85), Inches(1.9), Inches(1.3), BG_SLATE, accent)
    add_textbox(slide, x + Inches(0.1), Inches(5.9), Inches(1.7), Inches(0.3),
                title, font_size=12, color=accent, bold=True)
    add_textbox(slide, x + Inches(0.1), Inches(6.2), Inches(1.7), Inches(0.85),
                desc, font_size=10, color=GRAY_LIGHT)

# ========================================
# SLIDE 5 - PORTFOLIO & BLOG
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 5)
add_section_number(slide, 4)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Portfolio Interactivo & Blog Técnico", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), PURPLE)

# Portfolio section
add_card(slide, Inches(1), Inches(2), Inches(5.5), Inches(4.8), BG_CARD, PURPLE)
add_textbox(slide, Inches(1.2), Inches(2.15), Inches(5), Inches(0.4),
            "/portfolio/  —  CV Interactivo", font_size=20, color=PURPLE, bold=True)

portfolio_items = [
    "Animación de intro cinematográfica",
    "14+ posiciones con tags de tecnología y modales de detalle",
    "Proyectos GitHub filtrables por tecnología",
    "Stack tecnológico: Fullstack, ML, AI, Cloud, BI",
    "109+ certificaciones con categorización",
    "Educación timeline (Bachiller → Posgrado UAI)",
    "Dashboard económico comparador de inversiones 12 meses",
    "Newsletter 'The Data Digest' — suscripción semanal",
    "3 temas: Dark (default), Light, Terminal (CLI interactivo)",
    "Modo Terminal: comandos about, experience, neofetch, matrix",
    "Bilingüe completo (ES / EN) switcheable desde navbar",
    "Descarga de CV en PDF generado con jsPDF",
]
add_bullet_list(slide, Inches(1.2), Inches(2.6), Inches(5.1), Inches(4),
                portfolio_items, font_size=13, color=GRAY_LIGHT, bullet_char="▸")

# Blog section
add_card(slide, Inches(7), Inches(2), Inches(5.5), Inches(4.8), BG_CARD, CYAN)
add_textbox(slide, Inches(7.2), Inches(2.15), Inches(5), Inches(0.4),
            "/blog/  —  Data Engineering en las Trincheras", font_size=18, color=CYAN, bold=True)

blog_items = [
    "80+ artículos técnicos en Markdown",
    "Temas: BigQuery, dbt, Airflow, Docker, Kafka, embeddings",
    "Posts generados automáticamente por n8n desde Hacker News",
    "Publicación cruzada en dev.to con canonical URL",
    "Videos de YouTube embebidos (masterclasses)",
    "Filtro visual entre posts y videos",
    "Syntax highlighting + tablas + tipografía optimizada",
    "Fechas 2024-2026 con cobertura continua",
    "SEO optimizado con structured data por post",
    "Lead magnets contextuales en posts populares",
]
add_bullet_list(slide, Inches(7.2), Inches(2.6), Inches(5.1), Inches(4),
                blog_items, font_size=13, color=GRAY_LIGHT, bullet_char="▸")

# ========================================
# SLIDE 6 - HERRAMIENTAS FINANCIERAS
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 6)
add_section_number(slide, 5)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Herramientas Financieras — /recursos/", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), GREEN)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
            "Herramientas gratuitas para el contexto económico argentino — motor de SEO y generación de leads.",
            font_size=16, color=GRAY_LIGHT)

tools = [
    ("Calculadora de Ganancias", [
        "Art. 94 LIG — escala completa con deducciones Art. 30",
        "Valores AFIP 2026 actualizados",
        "Gráfico de torta de distribución impositiva",
        "Escenarios precargados (soltero, casado, hijos)",
        "Generación de reporte descargable",
    ], SKY),
    ("Cotizaciones del Dólar", [
        "6 tipos: oficial, blue, MEP, CCL, cripto, tarjeta",
        "Datos en tiempo real vía API",
        "Histórico interactivo con Chart.js",
        "Comparativa visual entre cotizaciones",
        "Frecuencia de uso: alta (tráfico recurrente)",
    ], GREEN),
    ("Indicadores Económicos", [
        "Inflación mensual e interanual",
        "Tasa de plazo fijo",
        "Riesgo país (EMBI Argentina)",
        "Históricos con gráficos interactivos",
        "Fuente de tráfico SEO de largo alcance",
    ], AMBER),
]

for i, (title, items, accent) in enumerate(tools):
    y = Inches(2.5) + Inches(i * 1.6)
    add_card(slide, Inches(1), y, Inches(3.8), Inches(1.4), BG_SLATE, accent)
    add_textbox(slide, Inches(1.15), y + Inches(0.05), Inches(3.5), Inches(0.3),
                title, font_size=15, color=accent, bold=True)
    add_bullet_list(slide, Inches(1.15), y + Inches(0.35), Inches(3.5), Inches(1),
                    items, font_size=11, color=GRAY_LIGHT, bullet_char="•", spacing=Pt(2))

add_card(slide, Inches(5.3), Inches(2.5), Inches(7.5), Inches(4.5), BG_CARD, GREEN)
add_textbox(slide, Inches(5.5), Inches(2.6), Inches(7), Inches(0.4),
            "Funciones Comerciales de las Herramientas", font_size=20, color=GREEN, bold=True)

commercial = [
    "SEO de largo alcance: capturan búsquedas informativas de alto volumen (calculadora ganancias 2026, dólar blue hoy, inflación argentina)",
    "Tráfico recurrente: el dólar y los indicadores económicos se consultan diariamente, generando visitas repetidas al dominio",
    "Autoridad de marca: posicionan a Mariano como referente que entiende tanto la tecnología como el contexto económico local",
    "Puerta de entrada: usuarios que llegan por las herramientas descubren los servicios de consultoría en la navegación",
    "Datos para lead scoring: el uso de las herramientas permite inferir perfil del visitante (profesional, empresario, inversor)",
    "Contenido compartible: herramientas que se comparten en redes y grupos de WhatsApp, generando backlinks orgánicos",
    "Diferenciador: muy pocos consultores tech ofrecen herramientas propias — genera valor percibido inmediato",
]
add_bullet_list(slide, Inches(5.5), Inches(3.1), Inches(7), Inches(3.5),
                commercial, font_size=14, color=GRAY_LIGHT, spacing=Pt(4))

# ========================================
# SLIDE 7 - SITIOS DE CLIENTES
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 7)
add_section_number(slide, 6)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Sitios de Clientes — Portfolio Comercial", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), TEAL)

# Neil
add_card(slide, Inches(1), Inches(2), Inches(5.5), Inches(5), BG_CARD, PURPLE)
add_textbox(slide, Inches(1.2), Inches(2.15), Inches(5), Inches(0.4),
            "Neil Climatizadores", font_size=22, color=PURPLE, bold=True)
add_textbox(slide, Inches(1.2), Inches(2.6), Inches(5), Inches(0.3),
            "www.mgatc.com/neil-site/  •  Climatización vehicular", font_size=13, color=GRAY_MID)
add_rect(slide, Inches(1.2), Inches(3), Inches(5), Inches(0.02), PURPLE)

neil_items = [
    "Empresa con patente exclusiva AR-031005B1 (Pre-Enfriado)",
    "Disponible en 6 idiomas: ES, EN, IT, FR, PT, DE",
    "Showcase de productos: climatizadores, AA, calderas, solar",
    "Sistema de Pre-Enfriamiento patentado",
    "Presencia internacional con representantes en Europa",
    "Formulario de contacto/cotización → webhook n8n",
    "basePath: /neil-site — integrado en el mismo deploy",
]
add_bullet_list(slide, Inches(1.2), Inches(3.2), Inches(5), Inches(3.5),
                neil_items, font_size=14, color=GRAY_LIGHT)

# El Portugues
add_card(slide, Inches(7), Inches(2), Inches(5.5), Inches(5), BG_CARD, TEAL)
add_textbox(slide, Inches(7.2), Inches(2.15), Inches(5), Inches(0.4),
            "El Portugués S.A.", font_size=22, color=TEAL, bold=True)
add_textbox(slide, Inches(7.2), Inches(2.6), Inches(5), Inches(0.3),
            "www.mgatc.com/elportugues-site/  •  Logística y distribución", font_size=13, color=GRAY_MID)
add_rect(slide, Inches(7.2), Inches(3), Inches(5), Inches(0.02), TEAL)

portugues_items = [
    "80+ años de historia — 4 generaciones (desde 1927)",
    "Certificación ISO 9001:2015 de calidad",
    "Servicios: transporte, distribución AMBA/interior",
    "Almacenamiento de 8.200 m²",
    "Propuesta comercial de automatización interna incluida",
    "Formulario de contacto → webhook n8n",
    "basePath: /elportugues-site — mismo pipeline de deploy",
]
add_bullet_list(slide, Inches(7.2), Inches(3.2), Inches(5), Inches(3.5),
                portugues_items, font_size=14, color=GRAY_LIGHT)

# ========================================
# SLIDE 8 - AUTOMATIZACIONES N8N
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 8)
add_section_number(slide, 7)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Automatizaciones — n8n Workflows", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), AMBER)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
            "6 workflows productivos auto-hospedados en mgobeaalcoba.app.n8n.cloud",
            font_size=16, color=GRAY_LIGHT)

workflows = [
    ("01", "AI Blog Creator", "Schedule semanal", "Lee Hacker News trends → genera artículo bilingüe con OpenRouter → publica en dev.to → commit a GitHub → registra en Supabase", SKY),
    ("02", "Webhook Contact Form", "Webhook", "Recibe formularios de neil-site y elportugues-site → log en Google Sheets → email de notificación + agradecimiento", GREEN),
    ("03", "Newsletter Subscription", "Webhook", "Captura suscripciones 'The Data Digest' → log en Google Sheets → email de bienvenida automatizado", PURPLE),
    ("04", "Free Automation Lead", "Webhook", "Leads del form 'Automatización Gratis' → log en Google Sheets → email de confirmación + notificación admin", AMBER),
    ("05", "Daily Content Digest", "Schedule diario", "Genera resumen de noticias tech vía AI + SerpAPI → email a suscriptores con curaduría automatizada", CYAN),
    ("06", "Unified AI Agent", "Webhook", "Agente RAG con búsqueda semántica en blog + captura leads en Supabase + calculadora de impuestos integrada", RED_ACCENT),
]

for i, (num, title, trigger, desc, accent) in enumerate(workflows):
    y = Inches(2.5) + Inches(i * 0.8)
    add_card(slide, Inches(1), y, Inches(11.3), Inches(0.7), BG_SLATE, accent)
    add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(0.5), Inches(0.3),
                num, font_size=16, color=accent, bold=True, font=FONT_CODE)
    add_textbox(slide, Inches(1.6), y + Inches(0.02), Inches(3), Inches(0.3),
                title, font_size=16, color=WHITE, bold=True)
    add_textbox(slide, Inches(4.6), y + Inches(0.05), Inches(1.5), Inches(0.25),
                trigger, font_size=11, color=accent, font=FONT_CODE)
    add_textbox(slide, Inches(4.6), y + Inches(0.3), Inches(7.5), Inches(0.35),
                desc, font_size=12, color=GRAY_LIGHT)

# ========================================
# SLIDE 9 - INFRAESTRUCTURA & DEPLOY
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 9)
add_section_number(slide, 8)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Infraestructura, Deploy & Base de Datos", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), SKY)

# Hosting
add_textbox(slide, Inches(1), Inches(2), Inches(5), Inches(0.4),
            "Hosting & Infraestructura", font_size=20, color=SKY, bold=True)

hosting_items = [
    "Cloudflare Pages — hosting principal, CDN global",
    "Cloudflare Registrar — dominio mgatc.com",
    "Cloudflare DNS — MX, CNAME, TXT, DMARC configurados",
    "Cloudflare Email Routing + Brevo SMTP — mariano@mgatc.com",
    "Security headers: X-Frame-Options, CSP, HSTS via _headers",
    "GitHub Pages — legacy redirect (gh-pages branch)",
]
add_bullet_list(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(2.5),
                hosting_items, font_size=14, color=GRAY_LIGHT)

# CI/CD
add_textbox(slide, Inches(1), Inches(4.5), Inches(5), Inches(0.4),
            "CI/CD Pipeline", font_size=20, color=GREEN, bold=True)

cicd_items = [
    "GitHub Actions — push a main → rebuild + redeploy",
    "Builds en paralelo para los 3 sitios (cv, neil, elportugues)",
    "Script root build-cloudflare.sh — ensamblado final",
    "Output: _site/ (cv root + neil-site + elportugues-site)",
    "Deploy automático a Cloudflare Pages sin intervención",
]
add_bullet_list(slide, Inches(1), Inches(5), Inches(5.5), Inches(2),
                cicd_items, font_size=14, color=GRAY_LIGHT)

# Database
add_textbox(slide, Inches(7), Inches(2), Inches(5.5), Inches(0.4),
            "Supabase — Base de Datos", font_size=20, color=PURPLE, bold=True)

db_items = [
    "PostgreSQL con Row Level Security (read-only anon)",
    "24+ tablas gestionadas · 11 migraciones SQL",
    "experience + experience_tags — 14+ posiciones",
    "projects + project_tags — GitHub projects filtrables",
    "education + certifications — 109+ certificaciones",
    "blog_posts — metadata de 80+ artículos",
    "consulting_packs + features + automations",
    "consulting_case_studies + tags",
    "Contenido dinámico: editar SQL = actualizar sin deploy",
]
add_bullet_list(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(4.5),
                db_items, font_size=14, color=GRAY_LIGHT)

# ========================================
# SLIDE 10 - SEO & ANALYTICS
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 10)
add_section_number(slide, 9)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "SEO, Analytics & Tracking", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), GREEN)

# SEO
add_textbox(slide, Inches(1), Inches(2), Inches(5), Inches(0.4),
            "SEO & Discoverabilidad", font_size=20, color=GREEN, bold=True)

seo_items = [
    "Sitemap dinámico (/sitemap.xml) — rutas + blog posts",
    "robots.txt con directivas por crawler (Google, Bing, AI bots)",
    "llms.txt para discoverabilidad por agentes de IA (GPTBot, Claude-Web)",
    "JSON-LD Schema.org: Person, WebSite, ProfessionalService",
    "Canonical URLs → https://www.mgatc.com",
    "Google Search Console verificado por DNS TXT",
    "Geo-meta tags para Buenos Aires, Argentina",
    "OpenGraph + Twitter Card metadata optimizados",
    "Redirect strategy: 301s a nivel CDN (Cloudflare _redirects)",
    "SEO plan documentado en 4 fases",
]
add_bullet_list(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(4),
                seo_items, font_size=13, color=GRAY_LIGHT, bullet_char="▸")

# Analytics
add_textbox(slide, Inches(7), Inches(2), Inches(5.5), Inches(0.4),
            "Analytics & Tracking", font_size=20, color=SKY, bold=True)

analytics_items = [
    "Google Analytics 4 — property G-DG0SLT5RY3 (compartida 3 sitios)",
    "Cross-domain tracking: mgatc.com, mgobeaalcoba.github.io, mgatc.pages.dev",
    "Cloudflare Web Analytics — cookieless, inmune a ad-blockers",
    "Key events: cv_download, sign_up, generate_lead, contact_click",
    "Eventos adicionales: contact_channel_select, whatsapp_click",
    "GA4 inline en <head> para mejor tag quality score",
    "ScrollTracker para medir engagement por sección",
    "Event tracking en calculadoras, herramientas y formularios",
    "Service view tracking por tipo de servicio",
    "Proposal modal open tracking",
]
add_bullet_list(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(4),
                analytics_items, font_size=13, color=GRAY_LIGHT, bullet_char="▸")

# ========================================
# SLIDE 11 - PRICING
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 11)
add_section_number(slide, 10)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Pricing — Modelo de Negocio", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), AMBER)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
            "Precios transparentes en USD — segmentados por B2B (empresas) y B2C (profesionales)",
            font_size=16, color=GRAY_LIGHT)

# B2B Column
add_card(slide, Inches(1), Inches(2.5), Inches(3.6), Inches(4.5), BG_CARD, SKY)
add_textbox(slide, Inches(1.2), Inches(2.6), Inches(3.2), Inches(0.4),
            "B2B — Empresas", font_size=18, color=SKY, bold=True)

b2b_rows = [
    ("Automatización (proyecto)", "$600"),
    ("Chatbot IA (proyecto)", "$800"),
    ("Dashboard BI (proyecto)", "$800"),
    ("Recurrente (soporte)", "$100-200/mes"),
    ("Primera automatización", "GRATIS"),
]
y = Inches(3.1)
for label, price in b2b_rows:
    add_textbox(slide, Inches(1.2), y, Inches(2.2), Inches(0.3),
                label, font_size=13, color=GRAY_LIGHT)
    add_textbox(slide, Inches(3.4), y, Inches(1), Inches(0.3),
                price, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.3)

add_textbox(slide, Inches(1.2), Inches(5.2), Inches(3.2), Inches(0.3),
            "ROI: 300-500% primer año", font_size=13, color=GREEN, bold=True)

# B2C Column
add_card(slide, Inches(5), Inches(2.5), Inches(3.6), Inches(4.5), BG_CARD, AMBER)
add_textbox(slide, Inches(5.2), Inches(2.6), Inches(3.2), Inches(0.4),
            "B2C — Profesionales", font_size=18, color=AMBER, bold=True)

b2c_rows = [
    ("Mentoría (sesión)", "$100"),
    ("Mentoría (pack 4)", "$350"),
    ("Mentoría (mensual)", "$280/mes"),
    ("Cursos intensivos", "$600-850"),
]
y = Inches(3.1)
for label, price in b2c_rows:
    add_textbox(slide, Inches(5.2), y, Inches(2.2), Inches(0.3),
                label, font_size=13, color=GRAY_LIGHT)
    add_textbox(slide, Inches(7.4), y, Inches(1), Inches(0.3),
                price, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.3)

add_textbox(slide, Inches(5.2), Inches(5.2), Inches(3.2), Inches(0.3),
            "Máx 4 personas por curso", font_size=13, color=AMBER, bold=True)

# Recruiting Column
add_card(slide, Inches(9), Inches(2.5), Inches(3.6), Inches(4.5), BG_CARD, TEAL)
add_textbox(slide, Inches(9.2), Inches(2.6), Inches(3.2), Inches(0.4),
            "Reclutamiento Tech", font_size=18, color=TEAL, bold=True)

rec_rows = [
    ("Success-based fee", "$3,000"),
    ("Retainer (1 búsqueda)", "$1,500/mes"),
    ("Retainer (2 búsquedas)", "$2,000/mes"),
    ("Garantía de reemplazo", "3 meses"),
]
y = Inches(3.1)
for label, price in rec_rows:
    add_textbox(slide, Inches(9.2), y, Inches(2.2), Inches(0.3),
                label, font_size=13, color=GRAY_LIGHT)
    add_textbox(slide, Inches(11.4), y, Inches(1), Inches(0.3),
                price, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.3)

add_textbox(slide, Inches(9.2), Inches(5.2), Inches(3.2), Inches(0.3),
            "Proceso: 3 semanas promedio", font_size=13, color=TEAL, bold=True)

# Benchmark note
add_card(slide, Inches(1), Inches(7.0), Inches(11.3), Inches(0.35), BG_SLATE)
add_textbox(slide, Inches(1.1), Inches(7.0), Inches(11), Inches(0.3),
            "Benchmark: 40-90% más económico que competencia directa. Posicionamiento: calidad premium a precio medio-competitivo.",
            font_size=12, color=GRAY_MID, bold=True)

# ========================================
# SLIDE 12 - LEAD GEN & CONVERSION
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 12)
add_section_number(slide, 11)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Lead Generation & Conversión", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), GREEN)

# Lead magnets
add_textbox(slide, Inches(1), Inches(2), Inches(5), Inches(0.4),
            "Embudo de Leads", font_size=20, color=GREEN, bold=True)

funnel = [
    "TOP: Herramientas financieras (SEO) — tráfico orgánico masivo",
    "MID: Blog técnico (80+ posts) — autoridad y confianza",
    "MID: Newsletter 'The Data Digest' — nurturing semanal",
    "BOTTOM: 'Automatización Gratis' — lead magnet principal",
    "BOTTOM: Contact form — WhatsApp directo + email",
    "BOTTOM: Proposal Modal — cotización personalizada en PDF",
    "BOTTOM: Floating CTA — siempre visible durante scroll",
]
add_bullet_list(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(3.5),
                funnel, font_size=14, color=GRAY_LIGHT, bullet_char="▸")

# Conversion tools
add_textbox(slide, Inches(7), Inches(2), Inches(5.5), Inches(0.4),
            "Herramientas de Conversión", font_size=20, color=SKY, bold=True)

conversion = [
    "WhatsApp directo (wa.me/5491127475569) — canal preferido en Argentina",
    "Email profesional mariano@mgatc.com — routing Cloudflare + Brevo SMTP",
    "Proposal Modal — generación de propuesta PDF con jsPDF",
    "Service modals — detalle de cada servicio con CTA contextual",
    "ScrollTracker — mide engagement por sección para optimizar",
    "GA4 key events — tracking completo del funnel de conversión",
    "n8n automation — respuesta automática en <24hs a leads",
    "Google Sheets backup — redundancia de datos de leads",
]
add_bullet_list(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(3.5),
                conversion, font_size=14, color=GRAY_LIGHT, bullet_char="▸")

# Funnel visual
add_card(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(1.3), BG_SLATE)
add_textbox(slide, Inches(1.1), Inches(5.85), Inches(11), Inches(0.3),
            "Funnel de Conversión Visual", font_size=16, color=WHITE, bold=True)

stages = [
    ("Tráfico Orgánico", "Herramientas + Blog + SEO"),
    ("Engagement", "Scroll + clicks + modales"),
    ("Captura", "Newsletter + Free Automation"),
    ("Conversión", "WhatsApp + Email + Proposal"),
]
for i, (stage, desc) in enumerate(stages):
    x = Inches(1.5) + Inches(i * 2.9)
    w = Inches(2.6)
    add_rect(slide, x, Inches(6.25), w, Inches(0.6), [SKY, PURPLE, AMBER, GREEN][i])
    add_textbox(slide, x, Inches(6.27), w, Inches(0.25),
                stage, font_size=13, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(6.52), w, Inches(0.25),
                desc, font_size=10, color=BG_DARK, alignment=PP_ALIGN.CENTER)

# ========================================
# SLIDE 13 - METRICAS & KPIs
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 13)
add_section_number(slide, 12)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Métricas & KPIs del Proyecto", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), CYAN)

# Quantitative
add_textbox(slide, Inches(1), Inches(2), Inches(5), Inches(0.4),
            "Métricas Cuantitativas", font_size=20, color=CYAN, bold=True)

metrics_q = [
    "3 sitios web en producción bajo un solo dominio",
    "80+ artículos técnicos publicados (blog)",
    "14+ experiencias profesionales documentadas",
    "109+ certificaciones registradas en Supabase",
    "24+ tablas de base de datos activas",
    "11 migraciones SQL versionadas",
    "6 workflows de automatización n8n productivos",
    "6 servicios de consultoría definidos y con pricing",
    "3 packs de pricing (Free, Pro, IA & Analytics)",
    "6 idiomas disponibles (Neil Climatizadores)",
    "3 temas UI (Dark, Light, Terminal)",
    "2 lenguajes completos (ES / EN)",
]
add_bullet_list(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(4.5),
                metrics_q, font_size=13, color=GRAY_LIGHT, bullet_char="▸")

# Qualitative
add_textbox(slide, Inches(7), Inches(2), Inches(5.5), Inches(0.4),
            "Logros Cualitativos", font_size=20, color=GREEN, bold=True)

metrics_qual = [
    "Monorepo autocontenido — todo en un solo repositorio",
    "CI/CD completamente automatizado (push → deploy)",
    "Contenido dinámico sin tocar código (Supabase + n8n)",
    "SEO avanzado con soporte para crawlers de IA",
    "Lead generation automatizada de principio a fin",
    "Blog auto-generado con curaduría AI de tendencias",
    "Herramientas financieras como motor de tráfico SEO",
    "Cross-posting automático en dev.to con canonical",
    "Email routing profesional con dominio propio",
    "Seguridad: security headers, RLS, sin secrets en repo",
    "Performance: static export, CDN global, lazy loading",
    "Accesibilidad: 3 temas, CLI terminal, bilingual",
]
add_bullet_list(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(4.5),
                metrics_qual, font_size=13, color=GRAY_LIGHT, bullet_char="▸")

# ========================================
# SLIDE 14 - OPORTUNIDADES DE MEJORA
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 14)
add_section_number(slide, 13)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Oportunidades de Mejora Comercial & Marketing", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), RED_ACCENT)

# Criticas
add_textbox(slide, Inches(1), Inches(2), Inches(5), Inches(0.4),
            "Críticas (alto impacto, bajo esfuerzo)", font_size=18, color=RED_ACCENT, bold=True)

criticas = [
    "Agregar métricas de impacto en el hero (+X procesos automatizados, $Y ahorrados)",
    "CTA principal diferenciado: botón sticky 'Agendar llamada gratis 15 min'",
    "Activar newsletter 'The Data Digest' en Mailchimp (ya configurado)",
    "Casos de éxito con formato problema → solución → resultado numérico",
]
add_bullet_list(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(2),
                criticas, font_size=13, color=GRAY_LIGHT, bullet_char="🔴")

# Estrategicas
add_textbox(slide, Inches(7), Inches(2), Inches(5.5), Inches(0.4),
            "Estratégicas (alto impacto, medio esfuerzo)", font_size=18, color=AMBER, bold=True)

estrategicas = [
    "Landing pages por servicio (/automatizacion, /ia, /bi) para SEO y ads",
    "Programa de referidos formal ('referí y ganá 1 sesión gratis')",
    "Posts 'insignia' escritos 100% por Mariano (experiencia real)",
    "Registro en directorios: Clutch, Google My Business, LinkedIn Services",
    "Video pitch de 60-90 seg embebido en hero",
]
add_bullet_list(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(2.5),
                estrategicas, font_size=13, color=GRAY_LIGHT, bullet_char="🟡")

# Crecimiento
add_textbox(slide, Inches(1), Inches(4.7), Inches(5), Inches(0.4),
            "Crecimiento (medio-alto impacto, mayor esfuerzo)", font_size=18, color=GREEN, bold=True)

crecimiento1 = [
    "Producto entry-level productizado ('Audit de datos en 48hs — $150')",
    "Estrategia sistemática en LinkedIn (3 posts/semana)",
    "Sección FAQ para resolver objeciones de venta",
    "Integrar CRM (HubSpot free, Pipedrive) para pipeline visual",
]
add_bullet_list(slide, Inches(1), Inches(5.2), Inches(5.5), Inches(2),
                crecimiento1, font_size=13, color=GRAY_LIGHT, bullet_char="🟢")

crecimiento2 = [
    "CTA contextual en herramientas financieras ('¿Querés automatizar esto?')",
    "Landing dedicada de cursos con syllabus, fechas y countdown",
    "Bundles de servicios (UpSkilling, Transform Team, Career Accelerator)",
    "Tier premium para anclar precios (cursos 1-on-1, executive search)",
]
add_bullet_list(slide, Inches(7), Inches(5.2), Inches(5.5), Inches(2),
                crecimiento2, font_size=13, color=GRAY_LIGHT, bullet_char="🟢")

# ========================================
# SLIDE 15 - PROYECCIONES
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 15)
add_section_number(slide, 14)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Proyección de Ingresos — Año 1", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), GREEN)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
            "Escenario realista con pricing optimizado (benchmark 2026)",
            font_size=16, color=GRAY_LIGHT)

# Q1-Q2
add_card(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(3), BG_CARD, SKY)
add_textbox(slide, Inches(1.2), Inches(2.6), Inches(5), Inches(0.4),
            "Q1-Q2 — Construcción de Portfolio", font_size=18, color=SKY, bold=True)

q1q2 = [
    ("5 Automatizaciones gratis → 2 conversiones", "$1,200"),
    ("3 Mentorías (pack 4)", "$1,050"),
    ("1 Curso SQL", "$700"),
    ("TOTAL Q1-Q2", "$2,950"),
]
y = Inches(3.1)
for label, price in q1q2:
    is_total = "TOTAL" in label
    add_textbox(slide, Inches(1.3), y, Inches(3.8), Inches(0.3),
                label, font_size=14, color=GREEN if is_total else GRAY_LIGHT, bold=is_total)
    add_textbox(slide, Inches(5.2), y, Inches(1.1), Inches(0.3),
                price, font_size=14, color=GREEN if is_total else WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.35)

# Q3-Q4
add_card(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(3), BG_CARD, PURPLE)
add_textbox(slide, Inches(7.2), Inches(2.6), Inches(5), Inches(0.4),
            "Q3-Q4 — Con Tracción y Testimonios", font_size=18, color=PURPLE, bold=True)

q3q4 = [
    ("3 Proyectos automatización", "$1,800"),
    ("2 Proyectos IA & Analytics", "$1,600"),
    ("1 Recruiting success fee", "$3,000"),
    ("5 Mentorías (pack 4)", "$1,750"),
    ("2 Cursos intensivos", "$1,550"),
    ("TOTAL Q3-Q4", "$9,700"),
]
y = Inches(3.1)
for label, price in q3q4:
    is_total = "TOTAL" in label
    add_textbox(slide, Inches(7.3), y, Inches(3.8), Inches(0.3),
                label, font_size=14, color=GREEN if is_total else GRAY_LIGHT, bold=is_total)
    add_textbox(slide, Inches(11.2), y, Inches(1.1), Inches(0.3),
                price, font_size=14, color=GREEN if is_total else WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.3)

# Total
add_card(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(1.3), BG_SLATE, GREEN)
add_textbox(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.4),
            "Proyección Anual Consolidada", font_size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

totals = [
    ("Consultoría B2B (proyectos)", "$6,100"),
    ("Mentorías B2C", "$2,800"),
    ("Cursos grupales", "$2,250"),
    ("Recurrentes (3 clientes × $100 × 6 meses)", "$1,800"),
    ("TOTAL AÑO 1", "$12,950"),
    ("UPSIDE (clientes grandes $200-500/mes)", "+$2,400 a $6,000"),
]
y = Inches(6.3)
for label, price in totals:
    is_total = "TOTAL" in label or "UPSIDE" in label
    add_textbox(slide, Inches(1.5), y, Inches(7), Inches(0.25),
                label, font_size=13, color=GREEN if is_total else GRAY_LIGHT, bold=is_total)
    add_textbox(slide, Inches(8.5), y, Inches(3.5), Inches(0.25),
                price, font_size=13, color=GREEN if is_total else WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.23)

# ========================================
# SLIDE 16 - ROADMAP
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide, 16)
add_section_number(slide, 15)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
            "Roadmap & Próximos Pasos", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.04), SKY)

phases = [
    ("FASE 1 — Inmediato", "Semana 1-2", [
        "Activar newsletter 'The Data Digest' en Mailchimp",
        "Agregar métricas de impacto al hero del consulting",
        "Implementar CTA sticky 'Agendar llamada gratuita'",
        "Crear casos de éxito con formato problema → solución → ROI",
        "Ajustar pricing recruiting a $3,000 (ya documentado)",
    ], RED_ACCENT),
    ("FASE 2 — Corto Plazo", "Mes 1-2", [
        "Crear landing pages individuales por servicio (SEO + ads)",
        "Activar programa de referidos formal",
        "Escribir 3-5 posts 'insignia' con experiencia real de Mariano",
        "Registrar en directorios (Clutch, Google My Business, LinkedIn Services)",
        "Grabar video pitch de 60-90 seg",
    ], AMBER),
    ("FASE 3 — Mediano Plazo", "Mes 2-3", [
        "Integrar CRM (HubSpot free) para pipeline y follow-ups",
        "Agregar FAQ section en página de consulting",
        "Crear producto entry-level ('Audit de datos en 48hs — $150')",
        "CTA contextual en herramientas financieras",
        "Estrategia de contenido LinkedIn (3 posts/semana)",
    ], GREEN),
    ("FASE 4 — Largo Plazo", "Mes 3-6", [
        "Landing dedicada de cursos con syllabus, fechas, countdown",
        "Implementar bundles de servicios (UpSkilling, Transform Team)",
        "Agregar tier premium (cursos 1-on-1, executive search)",
        "Ajustar pricing gradual (+20-40% basado en testimonios)",
        "Explorar mercado LATAM (no solo Argentina)",
    ], SKY),
]

for i, (phase, timeline, items, accent) in enumerate(phases):
    y = Inches(1.9) + Inches(i * 1.4)
    add_card(slide, Inches(1), y, Inches(11.3), Inches(1.25), BG_SLATE, accent)
    add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(3), Inches(0.3),
                phase, font_size=15, color=accent, bold=True)
    add_textbox(slide, Inches(4.2), y + Inches(0.05), Inches(2), Inches(0.3),
                timeline, font_size=13, color=GRAY_MID, font=FONT_CODE)

    item_text = "  •  " + "     •  ".join(items)
    add_textbox(slide, Inches(1.1), y + Inches(0.35), Inches(11), Inches(0.85),
                item_text, font_size=12, color=GRAY_LIGHT)

# ========================================
# SLIDE 17 - CLOSING
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(7.5), BG_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), SKY)
add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), SKY)
add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), SKY)
add_rect(slide, Inches(13.25), Inches(0), Inches(0.08), Inches(7.5), SKY)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
            "MGA TECH CONSULTING", font_size=42, color=SKY, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(4.5), Inches(2.5), Inches(4.3), Inches(0.03), SKY)

add_textbox(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.5),
            "Tecnología de clase mundial accesible para PyMEs argentinas", font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

summary_items = [
    "Ecosistema digital completo: consultoría + portfolio + blog + herramientas + sitios de clientes",
    "Arquitectura moderna: Next.js 14, Supabase, n8n, Cloudflare, CI/CD automatizado",
    "6 servicios con pricing competitivo y modelo freemium (primera automatización gratis)",
    "Lead generation automatizada de extremo a extremo con 6 workflows n8n",
    "SEO avanzado con soporte para crawlers de IA (GPTBot, Claude-Web, CCBot)",
    "Oportunidad de crecimiento: +55% ingresos con ajustes de pricing ya documentados",
]
add_bullet_list(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(2.5),
                summary_items, font_size=15, color=GRAY_LIGHT, spacing=Pt(6))

add_textbox(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
            "www.mgatc.com  •  mariano@mgatc.com  •  Buenos Aires, Argentina",
            font_size=16, color=GRAY_MID, bold=True, alignment=PP_ALIGN.CENTER)

# ========================================
# SAVE
# ========================================
output_path = os.path.expanduser("~/Documents/mgatc/MGA_Tech_Consulting_Presentacion.pdf")
# Save as PPTX first
pptx_path = os.path.expanduser("~/Documents/mgatc/MGA_Tech_Consulting_Presentacion.pptx")
prs.save(pptx_path)
print(f"Saved to {pptx_path}")
